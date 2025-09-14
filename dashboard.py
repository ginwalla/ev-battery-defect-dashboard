# Dashboard.py
# Streamlit portfolio-ready EV Battery Defect Detection dashboard
# Run: pip install streamlit pandas numpy matplotlib python-pptx pillow openpyxl
# Then: streamlit run Dashboard.py

import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import io
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

st.set_page_config(page_title="EV Battery Defect Detection", layout="wide")

# -----------------------
# Helper utilities
# -----------------------
MAX_FILE_MB = 100

def prediction_to_csv_download(prediction_response, filename="predictions.csv"):
    """
    Converts a prediction response (DataFrame or list of dicts) to CSV bytes and displays a Streamlit download button.
    """
    # If already a DataFrame, use it; else convert
    if isinstance(prediction_response, pd.DataFrame):
        df = prediction_response
    else:
        df = pd.DataFrame(prediction_response)
    csv_bytes = df.to_csv(index=False).encode('utf-8')
    st.download_button(
        label="Download Predictions as CSV",
        data=csv_bytes,
        file_name=filename,
        mime="text/csv"
    )

def validate_file_size(uploaded_file):
    if uploaded_file is None:
        return False, "No file provided."
    size_mb = uploaded_file.size / (1024 * 1024)
    if size_mb > MAX_FILE_MB:
        return False, f"File too large ({size_mb:.1f} MB). Max allowed is {MAX_FILE_MB} MB."
    return True, None

@st.cache_data(show_spinner=False)
def load_csv_buffer(buffer):
    try:
        return pd.read_csv(buffer)
    except Exception:
        buffer.seek(0)
        return pd.read_excel(buffer)

def save_uploaded_file(uploaded_file, folder="uploads"):
    os.makedirs(folder, exist_ok=True)
    path = os.path.join(folder, uploaded_file.name)
    with open(path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    return path

def create_pptx(summary_text, arch_image_path=None, metrics=None, out_path="DefectiveCell_Report.pptx"):
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Defective Cell — Summary Report"
    slide.placeholders[1].text = "Auto-generated from portfolio demo"

    # Summary + metrics slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = summary_text
    p.font.size = Pt(14)

    if metrics:
        # Add a small bullet list with metrics
        for k, v in metrics.items():
            p = tf.add_paragraph()
            p.text = f"{k}: {v}"
            p.level = 1

    # Architecture slide (image)
    if arch_image_path and os.path.exists(arch_image_path):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Architecture"
        left = Inches(0.7)
        top = Inches(1.6)
        slide.shapes.add_picture(arch_image_path, left, top, width=Inches(8))

    prs.save(out_path)
    return out_path

# Simple model stub (replace this with real model call)
def model_predict_row(row):
    """
    Demo model: returns probability of defect based on temperature and voltage heuristics.
    Replace with network call to your model server or local inference.
    """
    temp = row.get("temperature", 0)
    volt = row.get("voltage", 0)
    # heuristic demo: higher temp increases defect prob; low voltage slightly increases it
    prob = min(1.0, max(0.0, (temp - 25) / 40 + (3.6 - volt) * 0.2))
    label = "REJECT" if prob >= 0.5 else "PASS"
    # feature contributions (for a SHAP-like bar)
    expl = {"temperature": (temp - 30) / 10, "voltage": (3.7 - volt) * 2, "current": 0.0}
    return {"label": label, "prob": prob, "explanation": expl}

# -----------------------
# Layout: sidebar nav
# -----------------------
st.sidebar.title("Defective Cell — Portfolio")
page = st.sidebar.radio("Navigation", ["Home", "Data", "Model Playground", "Validation", "Export", "About"])

# -----------------------
# Home
# -----------------------
if page == "Home":
    st.title("⚡ EV Battery Defect Detection — Portfolio Demo")
    st.markdown(
        """
        This portfolio demo shows a lightweight pipeline: dataset upload → simple defect detection → model preview (placeholder)
        → validations → export to PPTX. Replace the `model_predict_row` function with your model server call (virtual environment).
        """
    )
    st.subheader("Quick sample")
    sample = pd.DataFrame(
        [
            {"Cell ID": "C-001", "Status": "PASS", "Remarks": "OK"},
            {"Cell ID": "C-002", "Status": "REJECT", "Remarks": "Seal fault"},
            {"Cell ID": "C-003", "Status": "PASS", "Remarks": "OK"},
        ]
    )
    st.table(sample)

# -----------------------
# Data upload & preview
# -----------------------
elif page == "Data":
    st.title("Data — Upload & Preview")
    st.write("Upload a CSV of cell readings or use the synthetic data generator below.")
    col1, col2 = st.columns([2,1])

    with col1:
        uploaded = st.file_uploader("Upload battery_data.csv", type=["csv","xlsx","xls"])
        if uploaded is not None:
            ok, msg = validate_file_size(uploaded)
            if not ok:
                st.error(msg)
                st.stop()
            try:
                df = load_csv_buffer(uploaded)
                st.success(f"Loaded `{uploaded.name}` — {len(df)} rows")
            except Exception as e:
                st.error(f"Could not parse file: {e}")
                st.stop()
        else:
            st.info("No file uploaded. Use the synthetic generator in the right panel or upload a CSV.")

    with col2:
        st.subheader("Synthetic Data Settings")
        size_choice = st.radio("Size", ("Small (20)", "Large (100)", "Custom"))
        if size_choice == "Small (20)":
            n_cells = 20
        elif size_choice == "Large (100)":
            n_cells = 100
        else:
            n_cells = st.slider("Number of synthetic cells", min_value=20, max_value=500, step=5, value=75)

        randomize = st.checkbox("Randomize each run", value=True)
        if not randomize:
            np.random.seed(42)

        if st.button("Generate synthetic"):
            df = pd.DataFrame({
                "cell_id": [f"C{i:03}" for i in range(1, n_cells + 1)],
                "temperature": np.random.randint(20, 55, n_cells),
                "voltage": np.round(np.random.uniform(3.0, 4.2, n_cells), 2),
                "current": np.round(np.random.uniform(0.5, 2.0, n_cells), 2),
            })
            st.session_state["df"] = df
            st.success(f"Synthetic dataset created ({n_cells} rows).")

    # Load df from session if present
    df = st.session_state.get("df") if "df" in st.session_state else (locals().get("df") if "df" in locals() else None)

    if df is not None:
        st.subheader("Dataset preview")
        st.dataframe(df.head(200))

        csv_full = df.to_csv(index=False).encode("utf-8")
        st.download_button(
            label="📥 Download Full Dataset (CSV)",
            data=csv_full,
            file_name="full_dataset.csv",
            mime="text/csv",
        )

# -----------------------
# Model Playground
# -----------------------
elif page == "Model Playground":
    st.title("Model Playground — Single item run")
    st.write("Run one sample through the model to preview prediction and explanation. Replace the stub with your model endpoint.")

    # pick a row (if dataset present) or create a synthetic sample
    df = st.session_state.get("df", None)
    if df is not None and not df.empty:
        chosen_idx = st.selectbox("Select row (from loaded dataset)", df.index.tolist())
        sample_row = df.loc[chosen_idx].to_dict()
    else:
        st.info("No dataset loaded — using a synthetic sample.")
        sample_row = {
            "cell_id": "sample_001",
            "temperature": st.slider("temperature", 20, 70, 42),
            "voltage": round(st.slider("voltage", 3.0, 4.2, 3.8), 2),
            "current": round(st.slider("current", 0.1, 3.0, 1.0), 2)
        }

    if st.button("Run model (demo)"):
        # TODO: replace this call with actual model server call (requests.post to your endpoint)
        res = model_predict_row(sample_row)
         # Vibe: show actual feature values (temp, voltage, current) for the selected row in a small table
        st.subheader("Input features")
        feat_df = pd.DataFrame([sample_row])
        st.table(feat_df)               
        st.metric(label="Prediction", value=res["label"], delta=f"{int(res['prob']*100)}%")
        st.write("Feature contributions (SHAP-like):")
        expl = res["explanation"]
        # normalize for display
        items = list(expl.items())
        keys = [k for k, _ in items]
        vals = [v for _, v in items]
        fig, ax = plt.subplots()
        ax.barh(keys, vals)
        ax.set_xlabel("Contribution")
        st.pyplot(fig)

# -----------------------
# Validation & Visuals
# -----------------------
elif page == "Validation":
    st.title("Data Validation & Visuals")
    df = st.session_state.get("df", None)
    if df is None:
        st.warning("Load or generate a dataset first on the Data page.")
        st.stop()

    st.subheader("Dataset health checks")
    missing = df.isnull().sum()
    if missing.any():
        st.warning("⚠️ Missing values found:")
        st.write(missing[missing > 0])
    else:
        st.success("✅ No missing values detected.")

    invalid_temp = df[df.get("temperature", pd.Series()) < 0]
    invalid_volt = df[df.get("voltage", pd.Series()) <= 0]
    invalid_curr = df[df.get("current", pd.Series()) < 0]

    if not invalid_temp.empty:
        st.error("Invalid temperature values detected:")
        st.dataframe(invalid_temp)
    if not invalid_volt.empty:
        st.error("Invalid voltage values detected:")
        st.dataframe(invalid_volt)
    if not invalid_curr.empty:
        st.error("Invalid current values detected:")
        st.dataframe(invalid_curr)
    if invalid_temp.empty and invalid_volt.empty and invalid_curr.empty:
        st.success("✅ All sensor readings are within valid ranges.")

    # Defect detection using a threshold widget (keeps your original idea)
    st.subheader("Defect detection (threshold)")
    threshold = st.slider("Defect threshold (°C)", 30, 70, 45)
    defective = df[df["temperature"] > threshold]
    percent = len(defective) / len(df) * 100 if len(df) > 0 else 0

    c1, c2, c3 = st.columns(3)
    c1.metric("Rows", len(df))
    c2.metric("Defective count", len(defective))
    c3.metric("Defect rate", f"{percent:.1f}%")

    st.subheader("Defective Cells")
    st.dataframe(defective)

    # Plot: bar chart of defective temperatures (more polished)
    if not defective.empty:
        fig, ax = plt.subplots(figsize=(8,3))
        ax.bar(defective["cell_id"], defective["temperature"])
        ax.set_ylabel("Temp (°C)")
        ax.set_xticklabels(defective["cell_id"], rotation=45, ha="right")
        st.pyplot(fig)
    else:
        st.info("No defective cells detected with current threshold.")

    # Distribution plot
    st.subheader("Temperature distribution")
    fig2, ax2 = plt.subplots()
    ax2.hist(df["temperature"], bins=15)
    ax2.set_xlabel("Temperature (°C)")
    st.pyplot(fig2)

    # Correlation heatmap of numeric columns (optimized)
    st.subheader("Correlation Heatmap (Numeric Columns)")
    numeric_df = df.select_dtypes(include=["number"])
    if not numeric_df.empty and numeric_df.shape[1] > 1:
        with st.spinner("Generating correlation heatmap..."):
            corr = numeric_df.corr()
            annot = corr.shape[0] <= 10
            fig3, ax3 = plt.subplots(figsize=(min(6, 0.5 * corr.shape[0] + 3), min(4, 0.5 * corr.shape[1] + 2)))
            sns.heatmap(corr, annot=annot, cmap="coolwarm", ax=ax3, cbar=True, square=True)
            ax3.set_title("Feature Correlation Matrix")
            st.pyplot(fig3)
    else:
        st.info("Not enough numeric columns for correlation heatmap.")

    # Download defective subset
    if not defective.empty:
        csv = defective.to_csv(index=False).encode("utf-8")
        st.download_button("📥 Download defective_cells.csv", data=csv, file_name="defective_cells.csv", mime="text/csv")

# -----------------------
# Export
# -----------------------
elif page == "Export":
    st.title("Export — Generate stakeholder PPTX")
    st.write("Generates a simple PPTX containing an executive summary and architecture image (if uploaded).")

    summary = st.text_area("Executive summary (will appear on PPTX)", value="This report summarizes a demo pipeline and sample predictions for defective cell detection.")
    st.markdown("**Architecture image (optional)** — upload PNG/JPG to include in the PPTX.")
    arch_file = st.file_uploader("Upload architecture image", type=["png","jpg","jpeg","svg"])
    arch_path = None
    if arch_file is not None:
        ok, msg = validate_file_size(arch_file)
        if not ok:
            st.error(msg)
        else:
            arch_path = save_uploaded_file(arch_file, folder="uploads")
            st.success(f"Saved architecture image to `{arch_path}`")

    # Add threshold slider to Export page to ensure 'threshold' is defined
    threshold = st.slider("Defect threshold (°C) for metrics", 30, 70, 45)

    if st.button("Generate PPTX"):
        # metrics included in PPTX
        df = st.session_state.get("df", None)
        metrics = {
            "Dataset rows": len(df) if df is not None else "N/A",
            "Defect threshold (°C)": threshold,
        }
        out = create_pptx(summary, arch_path, metrics=metrics, out_path="DefectiveCell_Report.pptx")
        st.success(f"PPTX generated: {out}")
        with open(out, "rb") as f:
            data = f.read()
        st.download_button("Download PPTX", data=data, file_name="DefectiveCell_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# -----------------------
# About
# -----------------------
elif page == "About":
    st.title("About / Next steps")
    st.write(
    )

   